# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(7.5)
prs.slide_height = Inches(13.333)

W = 7.5  # 슬라이드 폭
BG = RGBColor(0x1e, 0x29, 0x3b)
WHITE = RGBColor(0xff, 0xff, 0xff)
BLUE = RGBColor(0x4a, 0x6c, 0xf7)
GREEN = RGBColor(0x22, 0xc5, 0x5e)
GRAY = RGBColor(0x94, 0xa3, 0xb8)
RED = RGBColor(0xef, 0x44, 0x44)

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_para(tf, text, size=16, bold=False, color=WHITE, bullet=False):
    p = tf.add_paragraph()
    p.text = ("\u2022  " + text) if bullet else text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = Pt(8)

# ===== 표지 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.5, 3, 6.5, 1.5, "BlogMaster", 48, True, BLUE, PP_ALIGN.CENTER)
add_text(slide, 0.5, 5, 6.5, 1, "네이버 플레이스\n블로그 자동 포스팅", 26, False, WHITE, PP_ALIGN.CENTER)
add_text(slide, 0.5, 7, 6.5, 0.8, "사용 설명서", 20, False, GRAY, PP_ALIGN.CENTER)
add_text(slide, 0.5, 8, 6.5, 0.6, "2026년 4월", 16, False, GRAY, PP_ALIGN.CENTER)

# ===== 목차 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "목차", 36, True, BLUE)
tf = add_text(slide, 0.8, 2, 6, 10, "", 20)
for item in ["1.  프로그램 설치", "2.  로그인", "3.  메인 화면 구성",
             "4.  크롤링 (F5)", "5.  포스트 생성 (F7)", "6.  자동 포스팅 (F8)",
             "7.  프롬프트 편집", "8.  설정", "9.  단축키", "10. 주의사항"]:
    add_para(tf, item, 22, color=WHITE)

# ===== 1. 설치 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "1. 프로그램 설치", 36, True, BLUE)
tf = add_text(slide, 0.8, 2.5, 6, 8, "", 20)
add_para(tf, "NaverBlogAuto_Install.exe 실행", 22, color=WHITE, bullet=True)
add_para(tf, "", 10)
add_para(tf, "설치 경로 지정 → 설치 완료", 22, color=WHITE, bullet=True)
add_para(tf, "", 10)
add_para(tf, "바탕화면에 NaverBlogAuto\n아이콘 생성", 22, color=WHITE, bullet=True)
add_para(tf, "", 20)
add_para(tf, "※ Chrome 브라우저가\n   설치되어 있어야 합니다", 18, color=GRAY)

# ===== 2. 로그인 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "2. 로그인", 36, True, BLUE)
tf = add_text(slide, 0.8, 2.5, 6, 8, "", 20)
add_para(tf, "프로그램 실행 시\n로그인 창이 나타납니다", 22, color=WHITE)
add_para(tf, "", 14)
add_para(tf, "아이디와 비밀번호 입력 후\n로그인", 22, color=WHITE, bullet=True)
add_para(tf, "", 10)
add_para(tf, "사용 기간 내에만\n로그인 가능", 22, color=WHITE, bullet=True)
add_para(tf, "", 10)
add_para(tf, "아이디/비밀번호 저장 기능\n지원", 22, color=WHITE, bullet=True)

# ===== 3. 메인 화면 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "3. 메인 화면", 36, True, BLUE)

add_text(slide, 0.6, 2.2, 6, 0.6, "좌측 패널", 26, True, GREEN)
tf = add_text(slide, 0.8, 3, 6, 4, "", 18)
add_para(tf, "키워드 입력 (지역+업종)", 20, color=WHITE, bullet=True)
add_para(tf, "크롤링 시작 (F5) / 중단 (F6)", 20, color=WHITE, bullet=True)
add_para(tf, "포스트 생성 (F7)", 20, color=WHITE, bullet=True)
add_para(tf, "포스트 보기 및 포스팅 (F8)", 20, color=WHITE, bullet=True)
add_para(tf, "크롤링 결과보기", 20, color=WHITE, bullet=True)

add_text(slide, 0.6, 6.5, 6, 0.6, "우측 패널", 26, True, GREEN)
tf2 = add_text(slide, 0.8, 7.3, 6, 4, "", 18)
add_para(tf2, "계정 선택 (1~9번)", 20, color=WHITE, bullet=True)
add_para(tf2, "대시보드 (작성 수, API 사용량)", 20, color=WHITE, bullet=True)
add_para(tf2, "실행 로그 (실시간 진행 상황)", 20, color=WHITE, bullet=True)

add_text(slide, 0.6, 10, 6, 0.6, "상단 헤더", 26, True, GREEN)
tf3 = add_text(slide, 0.8, 10.8, 6, 2, "", 18)
add_para(tf3, "프롬프트: 업종별 프롬프트 편집", 20, color=WHITE, bullet=True)
add_para(tf3, "설정: API 키, 네이버 계정 등", 20, color=WHITE, bullet=True)

# ===== 4. 크롤링 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "4. 크롤링 (F5)", 36, True, BLUE)
tf = add_text(slide, 0.8, 2.5, 6, 8, "", 20)
add_para(tf, "네이버 플레이스에서\n업체 정보를 자동 수집합니다", 22, color=WHITE)
add_para(tf, "", 14)
add_para(tf, "키워드 입력:\n\"중랑구요양원\"\n\"강서구헬스장\" 형태", 20, color=WHITE, bullet=True)
add_para(tf, "", 10)
add_para(tf, "수집 개수 설정\n(기본 100개)", 20, color=WHITE, bullet=True)
add_para(tf, "", 10)
add_para(tf, "F5 또는 크롤링 시작\n버튼 클릭", 20, color=WHITE, bullet=True)
add_para(tf, "", 14)
add_para(tf, "수집 정보:\n업체명, 주소, 카테고리,\n근처역, 태그 등", 18, color=GRAY)

# ===== 5. 포스트 생성 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "5. 포스트 생성 (F7)", 36, True, BLUE)
tf = add_text(slide, 0.8, 2.5, 6, 5, "", 20)
add_para(tf, "크롤링된 업체 선택\n→ AI가 자동 생성", 22, color=WHITE)
add_para(tf, "", 14)
add_para(tf, "F7 클릭 → 업체 체크\n→ 포스트 생성", 20, color=WHITE, bullet=True)
add_para(tf, "", 8)
add_para(tf, "GPT / Gemini 선택 가능", 20, color=WHITE, bullet=True)
add_para(tf, "", 8)
add_para(tf, "본문 + 제목 + 태그\n자동 생성", 20, color=WHITE, bullet=True)
add_para(tf, "", 8)
add_para(tf, "이미지 Pixabay에서\n자동 다운로드 + 저장", 20, color=WHITE, bullet=True)

add_text(slide, 0.6, 8.5, 6, 0.6, "상태 표시", 26, True, GREEN)
tf2 = add_text(slide, 0.8, 9.3, 6, 3, "", 20)
add_para(tf2, "빨간색 ●  미생성", 22, color=RED)
add_para(tf2, "초록색 ●  생성 완료", 22, color=GREEN)
add_para(tf2, "[완]  포스팅 완료", 22, color=BLUE)

# ===== 6. 포스팅 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "6. 자동 포스팅 (F8)", 36, True, BLUE)
tf = add_text(slide, 0.8, 2.5, 6, 5, "", 20)
add_para(tf, "생성된 포스트를\n네이버 블로그에 자동 업로드", 22, color=WHITE)
add_para(tf, "", 14)
add_para(tf, "F8 클릭 → 항목 체크\n→ 자동 포스팅", 20, color=WHITE, bullet=True)
add_para(tf, "", 8)
add_para(tf, "발행 간격 설정\n(예: 2시간)", 20, color=WHITE, bullet=True)
add_para(tf, "", 8)
add_para(tf, "첫 포스트 즉시/대기\n옵션 선택 가능", 20, color=WHITE, bullet=True)

add_text(slide, 0.6, 8, 6, 0.6, "발행 방식", 26, True, GREEN)
tf2 = add_text(slide, 0.8, 8.8, 6, 3, "", 20)
add_para(tf2, "즉시 발행: 첫 번째 포스트", 20, color=WHITE, bullet=True)
add_para(tf2, "예약 발행: 2번째부터\n간격 적용", 20, color=WHITE, bullet=True)
add_para(tf2, "카테고리 + 태그 자동 입력", 20, color=WHITE, bullet=True)

# ===== 7. 프롬프트 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "7. 프롬프트 편집", 36, True, BLUE)
tf = add_text(slide, 0.8, 2.5, 6, 9, "", 20)
add_para(tf, "업종별 AI 프롬프트를\n자유롭게 커스텀", 22, color=WHITE)
add_para(tf, "", 14)
add_para(tf, "업종 추가/삭제 가능", 20, color=WHITE, bullet=True)
add_para(tf, "", 8)
add_para(tf, "AI 자동 생성:\n업종 페르소나에 맞는\n프롬프트 자동 작성", 20, color=WHITE, bullet=True)
add_para(tf, "", 8)
add_para(tf, "본문 / 제목 프롬프트 분리", 20, color=WHITE, bullet=True)
add_para(tf, "", 8)
add_para(tf, "제목 키워드 형식:\nXX동 / XX역 / XX구 선택", 20, color=WHITE, bullet=True)
add_para(tf, "", 8)
add_para(tf, "Pixabay 검색어 1/2/3 슬롯\n(한글 자동 번역)", 20, color=WHITE, bullet=True)

# ===== 8. 설정 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "8. 설정", 36, True, BLUE)

add_text(slide, 0.6, 2.2, 6, 0.6, "AI 설정", 26, True, GREEN)
tf = add_text(slide, 0.8, 3, 6, 2, "", 20)
add_para(tf, "AI 제공자: GPT / Gemini", 20, color=WHITE, bullet=True)
add_para(tf, "API 키 최대 3개 등록", 20, color=WHITE, bullet=True)

add_text(slide, 0.6, 5.5, 6, 0.6, "네이버 계정", 26, True, GREEN)
tf2 = add_text(slide, 0.8, 6.3, 6, 3, "", 20)
add_para(tf2, "계정 1~9번 등록 가능", 20, color=WHITE, bullet=True)
add_para(tf2, "블로그 ID / 네이버 ID\n/ PW / 카테고리", 20, color=WHITE, bullet=True)
add_para(tf2, "계정별 데이터 완전 분리", 20, color=WHITE, bullet=True)

add_text(slide, 0.6, 9, 6, 0.6, "Pixabay", 26, True, GREEN)
tf3 = add_text(slide, 0.8, 9.8, 6, 2, "", 20)
add_para(tf3, "API 키 등록 (이미지 검색용)", 20, color=WHITE, bullet=True)

# ===== 9. 단축키 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "9. 단축키", 36, True, BLUE)
for i, (k, d) in enumerate([("F5", "크롤링 시작"),
                              ("F6", "중단"),
                              ("F7", "포스트 생성"),
                              ("F8", "포스트 보기 및 포스팅")]):
    y = 3 + i * 2
    add_text(slide, 1, y, 2, 0.8, k, 40, True, BLUE, PP_ALIGN.CENTER)
    add_text(slide, 3.5, y + 0.1, 3.5, 0.6, d, 24, False, WHITE)

# ===== 10. 주의 =====
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
add_text(slide, 0.6, 0.8, 6, 0.8, "10. 주의사항", 36, True, RED)
tf = add_text(slide, 0.8, 2.5, 6, 9, "", 20)
add_para(tf, "Chrome 브라우저\n필수 설치", 22, color=WHITE, bullet=True)
add_para(tf, "", 12)
add_para(tf, "보안문자(캡차) 발생 시\n수동 해결", 22, color=WHITE, bullet=True)
add_para(tf, "", 12)
add_para(tf, "Chrome 업데이트 시\n호환 문제 발생 가능", 22, color=WHITE, bullet=True)
add_para(tf, "", 12)
add_para(tf, "API 키 외부 노출 금지", 22, color=WHITE, bullet=True)
add_para(tf, "", 12)
add_para(tf, "포스팅 간격을 너무 짧게\n하면 네이버 제재 가능", 22, color=WHITE, bullet=True)

out = "C:/Users/kidth/OneDrive/Desktop/BlogMaster_사용설명서.pptx"
prs.save(out)
print(f"완료: {out} ({os.path.getsize(out)//1024}KB)")
